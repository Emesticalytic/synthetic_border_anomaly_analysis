"""
UK Border Anomaly Detection - Presentation Generator
Converts markdown to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_presentation():
    """Generate PowerPoint presentation from stakeholder markdown"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "UK Border Security\nAnomaly Detection System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Machine Learning for Enhanced Border Security\nJanuary 2026"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "🎯 Executive Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Best Performing Model: Ensemble (93% ROC-AUC, 88% F1)"
    
    p = tf.add_paragraph()
    p.text = "✅ All performance targets exceeded"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "⚡ Real-time processing: 95ms average"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "💰 £5M annual operational savings"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🛡️ 34% improvement in threat detection"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🎖️ Production ready with GDPR compliance"
    p.level = 1
    
    # Slide 3: Model Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🏆 Model Performance Comparison"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🥇 Ensemble Model - 93% ROC-AUC, 88% F1"
    
    p = tf.add_paragraph()
    p.text = "Combines 4 models for maximum accuracy"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🥈 XGBoost - 92% ROC-AUC, 87% F1"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "2x faster for real-time processing"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🥉 Random Forest - 90% ROC-AUC, 85% F1"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Reliable backup/failover system"
    p.level = 1
    
    # Slide 4: Business Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "💼 Business Impact"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Operational Excellence"
    
    p = tf.add_paragraph()
    p.text = "30% reduction in manual workload"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "15x processing capacity increase"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Security Improvement"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "87% detection rate vs 65% manual"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4.2% false positive rate (target < 5%)"
    p.level = 1
    
    # Slide 5: Financial Returns
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "💰 Financial Returns & ROI"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Annual Savings: £5M"
    
    p = tf.add_paragraph()
    p.text = "Initial Investment: £950K"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5-Year ROI: 2,100%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Payback Period: 2.8 months"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Officer Satisfaction: 89% positive feedback"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "94% want continued use of the system"
    p.level = 1
    
    # Slide 6: Deployment Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🚀 Deployment Roadmap"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 2 (Months 4-6): Expansion"
    
    p = tf.add_paragraph()
    p.text = "Gatwick & Manchester airports"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "£150K infrastructure budget"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Train 500 additional officers"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 3 (Months 7-12): National Rollout"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "All 8 major UK airports"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "International data sharing integration"
    p.level = 1
    
    # Slide 7: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "✨ Key System Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Real-Time Detection: 95ms response time"
    
    p = tf.add_paragraph()
    p.text = "Explainable AI: SHAP analysis for transparency"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Interactive Dashboard: Live monitoring & insights"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "GDPR Compliant: Ethical AI with full compliance"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "24/7 Operation: No fatigue, consistent performance"
    p.level = 1
    
    # Slide 8: Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📋 Recommendations"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "✅ Approve immediate deployment"
    
    p = tf.add_paragraph()
    p.text = "Primary: Ensemble model for maximum accuracy"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Backup: XGBoost for high-speed scenarios"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Authorize Phase 2 expansion"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Budget: £150K for infrastructure"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Establish monitoring framework"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Monthly performance reviews"
    p.level = 1
    
    # Final slide: Contact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You\n\nQuestions?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Emem A. - Senior Data Scientist\nJanuary 2026"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(18)
    contact_para.font.color.rgb = RGBColor(100, 100, 100)
    contact_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "outputs/UK_Border_Anomaly_Detection_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created: {output_file}")
    return output_file

if __name__ == "__main__":
    print("🎬 Generating PowerPoint presentation...")
    create_presentation()
    print("✨ Done! Open the file in outputs/ folder")
